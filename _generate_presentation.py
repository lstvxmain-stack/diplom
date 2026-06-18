#!/usr/bin/env python3
"""Generate diploma defense presentation (.pptx) — 11 slides"""

import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# Colors
C_DARK  = RGBColor(0x2C, 0x3E, 0x50)
C_ACCENT = RGBColor(0x29, 0x80, 0xB9)
C_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_RED   = RGBColor(0xE7, 0x4C, 0x3C)
C_GREEN = RGBColor(0x27, 0xAE, 0x60)

screenshot_dir = os.path.join(os.path.dirname(__file__), 'data', 'screenshots')

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rect(slide, l, t, w, h, color):
    return add_shape(slide, l, t, w, h, color)

def add_text_box(slide, left, top, w, h):
    return slide.shapes.add_textbox(left, top, w, h)

def tb_text(tb, text, size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT):
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = 'Calibri'
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf

def add_bullets(tb, items, size=16, color=C_DARK, spacing=Pt(8)):
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.level = 0
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = 'Calibri'
        run.font.size = Pt(size)
        run.font.color.rgb = color

def add_table(slide, left, top, w, h, headers, rows):
    tbl_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, w, h)
    tbl = tbl_shape.table
    for i, hdr in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(11); r.font.color.rgb = C_WHITE
                r.font.name = 'Calibri'
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_ACCENT
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(10); r.font.name = 'Calibri'
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_LIGHT
    return tbl_shape

# ============================================================
# SLIDE 1 — Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, C_DARK)

add_rect(slide, 0, 0, Inches(0.4), H, C_ACCENT)

tb_text(add_text_box(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(1)),
        'ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА', 36, True, C_WHITE, PP_ALIGN.LEFT)

tb_text(add_text_box(slide, Inches(1.2), Inches(2.8), Inches(11), Inches(1.5)),
        'Разработка информационно-справочной системы\nкультурно-развлекательных центров\nБелгородской области',
        28, False, C_LIGHT, PP_ALIGN.LEFT)

tb_text(add_text_box(slide, Inches(1.2), Inches(5.2), Inches(5), Inches(0.6)),
        'Обучающийся: [ФИО]', 18, False, C_LIGHT)
tb_text(add_text_box(slide, Inches(1.2), Inches(5.8), Inches(5), Inches(0.6)),
        'Руководитель: [ФИО, должность]', 18, False, C_LIGHT)

tb_text(add_text_box(slide, Inches(1.2), Inches(6.6), Inches(5), Inches(0.6)),
        'Белгород, 2026', 16, False, RGBColor(0x95, 0xA5, 0xA6))

# ============================================================
# SLIDE 2 — Актуальность
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_WHITE)
add_rect(slide, 0, 0, W, Inches(1.2), C_ACCENT)
tb_text(add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8)),
        'Актуальность и проблема', 30, True, C_WHITE, PP_ALIGN.LEFT)

add_bullets(add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.2)),
    ['1503 учреждения культуры в Белгородской области',
     'Информация разрознена по множеству ресурсов:',
     '  • bel.cultreg.ru — реестр, но нет карты и афиши',
     '  • afishka31.ru — афиша, но нет каталога',
     '  • Яндекс.Карты — карта, но нет спец.категоризации',
     '  • 2ГИС — справочник, нет афиши',
     'Гос. политика: нацпроект «Культура» — цифровизация',
     'Потребность жителей в едином сервисе'], 17, C_DARK)

tb_text(add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(1)),
        'Проблема', 22, True, C_RED, PP_ALIGN.LEFT)
tb_text(add_text_box(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(2.5)),
        'Ни один существующий ресурс не объединяет '
        'на одной платформе интерактивную карту, каталог '
        'учреждений, афишу мероприятий, административную '
        'панель и открытый API.', 18, False, C_DARK)

# Highlight box
box = add_shape(slide, Inches(7.0), Inches(4.8), Inches(5.5), Inches(1.5), RGBColor(0xFD, 0xED, 0xEC))
tb = add_text_box(slide, Inches(7.3), Inches(5.0), Inches(5.0), Inches(1.2))
tb_text(tb, '→ Разработка собственной ИТ-системы,\n   объединяющей все функции', 18, True, C_RED, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 3 — Цель и задачи
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_WHITE)
add_rect(slide, 0, 0, W, Inches(1.2), C_ACCENT)
tb_text(add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8)),
        'Цель и задачи', 30, True, C_WHITE, PP_ALIGN.LEFT)

tb_text(add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.6)),
        'Цель — разработка веб-приложения с интерактивной картой, каталогом и афишей', 20, True, C_ACCENT)

add_bullets(add_text_box(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4.5)),
    ['Анализ сферы и источников данных',
     'Обзор существующих решений-аналогов',
     'Разработка ТЗ (ГОСТ 19.201-78, ГОСТ 34.602-89)',
     'Обоснование технологического стека',
     'Реализация приложения',
     'Тестирование и оценка'], 18, C_DARK)

# Right column — numbers
tb = add_text_box(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(4.5))
tf = tb.text_frame; tf.word_wrap = True
steps = [
    ('1', 'Анализ', 'Выявлено 10 категорий,\n1503 учреждения'),
    ('2', 'Обзор', '6 аналогов, таблица\nпо 11 критериям'),
    ('3', 'ТЗ', '10 разделов по ГОСТ'),
    ('4', 'Технологии', 'Python, Flask, Leaflet,\nSQLite, Bootstrap, Render'),
    ('5', 'Реализация', 'Карта, каталог, афиша,\nадминка, API, парсеры'),
    ('6', 'Тестирование', '20 сценариев — все пройдены'),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(2.5) + Inches(0.75) * i
    # number circle
    s = add_shape(slide, Inches(7.0), y, Inches(0.5), Inches(0.5), C_ACCENT)
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = s.text_frame.paragraphs[0].add_run()
    r.text = num; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = C_WHITE; r.font.name = 'Calibri'
    # title
    t1 = add_text_box(slide, Inches(7.7), y - Pt(2), Inches(4.5), Inches(0.3))
    tb_text(t1, f'{title}:', 16, True, C_DARK)
    t2 = add_text_box(slide, Inches(7.7), y + Inches(0.25), Inches(4.5), Inches(0.4))
    tb_text(t2, desc, 14, False, RGBColor(0x7F, 0x8C, 0x8D))

# ============================================================
# SLIDE 4 — Анализ аналогов
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_WHITE)
add_rect(slide, 0, 0, W, Inches(1.2), C_ACCENT)
tb_text(add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8)),
        'Сравнительный анализ существующих решений', 30, True, C_WHITE, PP_ALIGN.LEFT)

headers = ['Критерий', 'Яндекс.\nКарты', '2ГИС', 'bel.\ncultreg.ru', 'Afishka\n31', 'Culture\n.РФ', 'Наша\nсистема']
rows = [
    ['Интерактивная карта', '✓', '✓', '—', '—', '—', '✓'],
    ['Кластеризация', '✓', '✓', '—', '—', '—', '✓'],
    ['Фильтр по категориям', '±', '±', '✓', '—', '✓', '✓'],
    ['Фильтр по районам', '—', '—', '✓', '—', '—', '✓'],
    ['Каталог учреждений', '✓', '✓', '✓', '—', '✓', '✓'],
    ['Афиша мероприятий', '—', '—', '—', '✓', '✓', '✓'],
    ['События на карте', '—', '—', '—', '—', '—', '✓'],
    ['Админ-панель', '—', '—', '—', '—', '—', '✓'],
    ['Открытый API', '✓', '✓', '✓', '—', '✓', '✓'],
    ['Открытый код', '—', '—', '—', '—', '—', '✓'],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5), headers, rows)

# ============================================================
# SLIDE 5 — Технологический стек
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_WHITE)
add_rect(slide, 0, 0, W, Inches(1.2), C_ACCENT)
tb_text(add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8)),
        'Технологический стек', 30, True, C_WHITE, PP_ALIGN.LEFT)

techs = [
    ('Python 3.12', 'Язык программирования', C_ACCENT),
    ('Flask 3.0', 'Веб-фреймворк', C_GREEN),
    ('SQLite + SQLAlchemy', 'База данных и ORM', C_RED),
    ('Leaflet + MarkerCluster', 'Интерактивная карта', RGBColor(0xE6, 0x7E, 0x22)),
    ('Bootstrap 5.3', 'CSS-фреймворк', RGBColor(0x8E, 0x44, 0xAD)),
    ('Flask-Admin', 'Админ-панель', C_DARK),
    ('Gunicorn', 'WSGI-сервер', RGBColor(0x2E, 0x86, 0xC1)),
    ('Render.com', 'Хостинг и деплой', RGBColor(0x16, 0xA0, 0x85)),
]
for i, (name, desc, color) in enumerate(techs):
    col = i % 4
    row = i // 4
    l = Inches(0.8) + Inches(3.1) * col
    t = Inches(1.6) + Inches(2.6) * row
    # Card background
    card = add_shape(slide, l, t, Inches(2.8), Inches(2.1), color)
    # Name
    tb1 = add_text_box(slide, l + Inches(0.2), t + Inches(0.3), Inches(2.4), Inches(0.8))
    tb_text(tb1, name, 20, True, C_WHITE, PP_ALIGN.CENTER)
    # Description
    tb2 = add_text_box(slide, l + Inches(0.2), t + Inches(1.2), Inches(2.4), Inches(0.6))
    tb_text(tb2, desc, 14, False, C_LIGHT, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6 — Архитектура
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_WHITE)
add_rect(slide, 0, 0, W, Inches(1.2), C_ACCENT)
tb_text(add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8)),
        'Архитектура системы', 30, True, C_WHITE, PP_ALIGN.LEFT)

# Client box
c1 = add_shape(slide, Inches(1.5), Inches(1.8), Inches(3.0), Inches(1.8), C_ACCENT)
tb_text(add_text_box(slide, Inches(1.7), Inches(2.0), Inches(2.6), Inches(0.5)),
        'Клиент (Браузер)', 20, True, C_WHITE, PP_ALIGN.CENTER)
tb_text(add_text_box(slide, Inches(1.7), Inches(2.6), Inches(2.6), Inches(0.8)),
        'HTML5 + CSS3 + JavaScript\nLeaflet, Bootstrap\nJinja2-шаблоны', 14, False, C_LIGHT, PP_ALIGN.CENTER)

# Arrow 1
tb_text(add_text_box(slide, Inches(4.7), Inches(2.3), Inches(1.5), Inches(0.5)),
        'HTTP →', 22, True, C_ACCENT, PP_ALIGN.CENTER)

# Server box
c2 = add_shape(slide, Inches(6.0), Inches(1.8), Inches(3.0), Inches(1.8), C_GREEN)
tb_text(add_text_box(slide, Inches(6.2), Inches(2.0), Inches(2.6), Inches(0.5)),
        'Сервер (Flask)', 20, True, C_WHITE, PP_ALIGN.CENTER)
tb_text(add_text_box(slide, Inches(6.2), Inches(2.6), Inches(2.6), Inches(0.8)),
        'Gunicorn WSGI\nSQLAlchemy ORM\nJSON API', 14, False, C_LIGHT, PP_ALIGN.CENTER)

# Arrow 2
tb_text(add_text_box(slide, Inches(9.2), Inches(2.3), Inches(1.5), Inches(0.5)),
        '↔ SQL', 22, True, C_GREEN, PP_ALIGN.CENTER)

# DB box
c3 = add_shape(slide, Inches(10.5), Inches(1.8), Inches(2.5), Inches(1.8), C_RED)
tb_text(add_text_box(slide, Inches(10.6), Inches(2.0), Inches(2.2), Inches(0.5)),
        'SQLite', 20, True, C_WHITE, PP_ALIGN.CENTER)
tb_text(add_text_box(slide, Inches(10.6), Inches(2.6), Inches(2.2), Inches(0.6)),
        '4 таблицы\n1503 записи', 14, False, C_LIGHT, PP_ALIGN.CENTER)

# Models section
tb_text(add_text_box(slide, Inches(0.8), Inches(4.0), Inches(5), Inches(0.5)),
        'Модели данных:', 20, True, C_DARK)

models = [
    ('Category (10)', 'id, name, icon, color, sort_order'),
    ('Venue (1503)', 'id, name, address, lat, lon, district,\nphone, website, description, category_id'),
    ('Event (10+)', 'id, title, description, date_start,\ndate_end, price, age_rating, venue_id'),
    ('AdminUser (1)', 'id, username, password_hash'),
]
for i, (name, fields) in enumerate(models):
    col = i % 2
    row = i // 2
    l = Inches(0.8) + Inches(6.0) * col
    t = Inches(4.5) + Inches(1.3) * row
    mbox = add_shape(slide, l, t, Inches(5.5), Inches(1.1), C_LIGHT)
    tb_text(add_text_box(slide, l + Inches(0.2), t + Inches(0.1), Inches(5.0), Inches(0.4)),
            name, 16, True, C_ACCENT)
    tb_text(add_text_box(slide, l + Inches(0.2), t + Inches(0.5), Inches(5.0), Inches(0.5)),
            fields, 12, False, C_DARK)

# API section
tb_text(add_text_box(slide, Inches(0.8), Inches(6.0), Inches(5), Inches(0.5)),
        'API-маршруты:', 20, True, C_DARK)
tb_text(add_text_box(slide, Inches(0.8), Inches(6.4), Inches(12), Inches(0.8)),
        'GET /api/map-data — GET /api/categories — GET /api/districts — GET / — GET /venues/ — GET /events/', 14, False, RGBColor(0x7F, 0x8C, 0x8D))

# ============================================================
# SLIDE 7 — Процесс разработки
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_WHITE)
add_rect(slide, 0, 0, W, Inches(1.2), C_ACCENT)
tb_text(add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8)),
        'Ключевые моменты разработки', 30, True, C_WHITE, PP_ALIGN.LEFT)

# Left column
tb_text(add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5)),
        'Импорт данных (cultreg.ru)', 20, True, C_ACCENT)
add_bullets(add_text_box(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(1.5)),
    ['REST API, пагинация (100/запрос)',
     'Маппинг 12 категорий → 10',
     'Импортировано 1487 объектов'], 15, C_DARK)

tb_text(add_text_box(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.5)),
        'Карта и кластеризация', 20, True, C_ACCENT)
add_bullets(add_text_box(slide, Inches(0.8), Inches(4.4), Inches(5.5), Inches(1.5)),
    ['1503 маркера через L.divIcon (цвет по категории)',
     'MarkerCluster (chunkedLoading, maxRadius=50)',
     'Граница области: GeoJSON (890 точек, красный контур)',
     'Popup с названием, адресом и 3 ближайшими событиями'], 15, C_DARK)

tb_text(add_text_box(slide, Inches(0.8), Inches(6.0), Inches(5.5), Inches(0.5)),
        'Фильтрация и поиск', 20, True, C_ACCENT)
add_bullets(add_text_box(slide, Inches(0.8), Inches(6.4), Inches(5.5), Inches(0.8)),
    ['По категории, району, названию',
     'Динамическая перезагрузка маркеров'], 15, C_DARK)

# Right column
tb_text(add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5)),
        'Оптимизация N+1 запроса', 20, True, C_RED)
box = add_shape(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(1.8), RGBColor(0xFD, 0xED, 0xEC))
tb = add_text_box(slide, Inches(7.3), Inches(2.4), Inches(5.0), Inches(1.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run(); r.text = '\u0414\u041e: 1504 SQL-\u0437\u0430\u043f\u0440\u043e\u0441\u043e\u0432, 12 \u0441\u0435\u043a\n'
r.font.size = Pt(14); r.font.name = 'Calibri'; r.font.color.rgb = C_RED
p2 = tf.add_paragraph(); r2 = p2.add_run()
r2.text = '\u041f\u041e\u0421\u041b\u0415: 2 SQL-\u0437\u0430\u043f\u0440\u043e\u0441\u0430, 5-6 \u0441\u0435\u043a'
r2.font.size = Pt(14); r.font.name = 'Calibri'; r.font.color.rgb = C_GREEN; r2.font.bold = True

tb_text(add_text_box(slide, Inches(7.0), Inches(4.3), Inches(5.5), Inches(0.5)),
        'Административная панель', 20, True, C_ACCENT)
add_bullets(add_text_box(slide, Inches(7.0), Inches(4.8), Inches(5.5), Inches(1.0)),
    ['Flask-Admin + Flask-Login',
     'Управление: категории, объекты, мероприятия',
     'Запуск парсеров из веб-интерфейса'], 15, C_DARK)

tb_text(add_text_box(slide, Inches(7.0), Inches(5.8), Inches(5.5), Inches(0.5)),
        'Деплой', 20, True, C_ACCENT)
add_bullets(add_text_box(slide, Inches(7.0), Inches(6.2), Inches(5.5), Inches(1.0)),
    ['Render.com + GitHub (автодеплой)',
     'Gunicorn + Procfile',
     'https://kulturnaya-karta-31.onrender.com'], 15, C_DARK)

# ============================================================
# SLIDE 8 — Скриншоты (4 коллажом)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_WHITE)
add_rect(slide, 0, 0, W, Inches(1.2), C_ACCENT)
tb_text(add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8)),
        'Интерфейс системы', 30, True, C_WHITE, PP_ALIGN.LEFT)

screenshots = [
    ('fig_a1_main.png', 'Главная страница — карта'),
    ('fig_a3_venue_detail.png', 'Детальная страница объекта'),
    ('fig_a4_events.png', 'Афиша мероприятий'),
    ('fig_a5_popup.png', 'Popup маркера'),
]
positions = [(Inches(0.5), Inches(1.5)), (Inches(6.7), Inches(1.5)),
             (Inches(0.5), Inches(4.5)), (Inches(6.7), Inches(4.5))]
for (fname, caption), (l, t) in zip(screenshots, positions):
    fpath = os.path.join(screenshot_dir, fname)
    if os.path.exists(fpath):
        slide.shapes.add_picture(fpath, l, t, Inches(5.8), Inches(2.7))
    tb_text(add_text_box(slide, l, t + Inches(2.7), Inches(5.8), Inches(0.4)),
            caption, 13, False, C_DARK, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9 — Тестирование
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_WHITE)
add_rect(slide, 0, 0, W, Inches(1.2), C_ACCENT)
tb_text(add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8)),
        'Тестирование', 30, True, C_WHITE, PP_ALIGN.LEFT)

add_bullets(add_text_box(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(2.5)),
    ['Метод: чёрный ящик (ручное тестирование)',
     '20 тестовых сценариев — все модули системы',
     'Функциональное тестирование API (4 сценария)',
     'Проверка кластеризации, фильтров, админки',
     'Результат: все 20 тестов пройдены'], 17, C_DARK)

# Performance box
perf_box = add_shape(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.2), C_LIGHT)
tb_text(add_text_box(slide, Inches(1.1), Inches(4.7), Inches(5.0), Inches(0.5)),
        'Производительность', 18, True, C_ACCENT)
tb_text(add_text_box(slide, Inches(1.1), Inches(5.3), Inches(5.0), Inches(1.2)),
        '• API /api/map-data (без фильтров): 5-6 сек\n'
        '• API /api/map-data (с фильтром): 0.1-0.5 сек\n'
        '• Полная загрузка страницы: 7-8 сек', 15, False, C_DARK)

# Results box
res_box = add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.1), RGBColor(0xE8, 0xF8, 0xF5))
tb_text(add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5.0), Inches(0.5)),
        'Итоги тестирования', 18, True, C_GREEN)
tb_text(add_text_box(slide, Inches(7.3), Inches(2.4), Inches(5.0), Inches(4.0)),
    '✓ Все функциональные требования ТЗ реализованы\n\n'
    '✓ Система работает штатно\n\n'
    '✓ Дефектов не выявлено\n\n'
    '✓ Система готова к эксплуатации', 16, False, C_DARK)

# ============================================================
# SLIDE 10 — Результаты
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_WHITE)
add_rect(slide, 0, 0, W, Inches(1.2), C_ACCENT)
tb_text(add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8)),
        'Результаты работы', 30, True, C_WHITE, PP_ALIGN.LEFT)

stats = [
    ('1503', 'учреждения\nв базе данных'),
    ('10', 'категорий\nкультуры'),
    ('20+', 'районов\nобласти'),
    ('6', 'страниц /\nмаршрутов'),
    ('11', 'критериев\nпревосходства'),
    ('27', 'источников\nв списке лит-ры'),
]
for i, (num, desc) in enumerate(stats):
    l = Inches(0.5) + Inches(2.1) * i
    t = Inches(1.6)
    card = add_shape(slide, l, t, Inches(1.9), Inches(1.8), C_ACCENT if i < 3 else C_GREEN)
    tb_text(add_text_box(slide, l + Inches(0.1), t + Inches(0.2), Inches(1.7), Inches(0.7)),
            num, 30, True, C_WHITE, PP_ALIGN.CENTER)
    tb_text(add_text_box(slide, l + Inches(0.1), t + Inches(1.0), Inches(1.7), Inches(0.6)),
            desc, 12, False, C_LIGHT, PP_ALIGN.CENTER)

# URL box
url_box = add_shape(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(1.0), C_LIGHT)
tb_text(add_text_box(slide, Inches(0.7), Inches(3.95), Inches(12.0), Inches(0.6)),
        'Веб-приложение: https://kulturnaya-karta-31.onrender.com     |     GitHub: github.com/lstvxmain-stack/diplom',
        18, True, C_ACCENT, PP_ALIGN.CENTER)

add_bullets(add_text_box(slide, Inches(0.8), Inches(5.2), Inches(5.5), Inches(2.0)),
    ['Практическая значимость:',
     '  • Жители и гости региона → планирование досуга',
     '  • Органы управления → мониторинг и анализ',
     '  • Разработчики → открытый JSON API'], 16, C_DARK)

add_bullets(add_text_box(slide, Inches(7.0), Inches(5.2), Inches(5.5), Inches(2.0)),
    ['Перспективы:',
     '  • Собственный домен',
     '  • PostgreSQL для производительности',
     '  • Мобильная версия',
     '  • Отзывы и рейтинги учреждений'], 16, C_DARK)

# ============================================================
# SLIDE 11 — Спасибо
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_DARK)
add_rect(slide, 0, 0, Inches(0.4), H, C_ACCENT)

tb_text(add_text_box(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(1.5)),
        'Спасибо за внимание!', 44, True, C_WHITE, PP_ALIGN.LEFT)
tb_text(add_text_box(slide, Inches(1.2), Inches(3.5), Inches(11), Inches(1)),
        'Готов(а) ответить на ваши вопросы', 24, False, C_LIGHT, PP_ALIGN.LEFT)

tb_text(add_text_box(slide, Inches(1.2), Inches(5.5), Inches(11), Inches(1.2)),
        'Репозиторий: github.com/lstvxmain-stack/diplom\n'
        'Сайт: https://kulturnaya-karta-31.onrender.com',
        18, False, RGBColor(0x95, 0xA5, 0xA6))

# SAVE
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Presentation_KK31.pptx')
prs.save(out)
print(f'Saved: {out}')
print(f'Size: {os.path.getsize(out)} bytes')
